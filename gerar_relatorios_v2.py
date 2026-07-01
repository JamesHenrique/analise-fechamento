"""
Gerar Relatórios com Identidade Visual da Logo
PPTX + PDF + HTML com logo e cores da marca Relive
"""
import json
from datetime import datetime
from io import BytesIO
from PIL import Image

BASE = 'G:/Outros computadores/NOTEBOOK1/James/projeto_relive/Relatório/Relatório de abertura-fechamento'
LOGO_PATH = f'{BASE}/logo.png'

# Carregar KPIs
with open(f'{BASE}/kpis.json', 'r', encoding='utf-8') as f:
    kpis = json.load(f)

# ============================================================
# PALETA DE CORES — Extraída da logo
# ============================================================
PRIMARY = '#602060'      # plum escuro
SECONDARY = '#a02060'    # berry/wine
ACCENT = '#e06020'       # coral/terracota
DARK_BG = '#420c38'      # quase preto plum
LIGHT_BG = '#f8f4fc'     # lavanda suave
TEXT = '#3c2850'         # texto escuro
WHITE = '#ffffff'
MUTED = '#8e7ca3'        # texto secundário
SUCCESS_COLOR = '#2e7d32'
WARNING_COLOR = '#e67e22'
DANGER_COLOR = '#c0392b'

def hex_to_rgb(h):
    h = h.lstrip('#')
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

def fmt_num(n):
    return f"{n:,}".replace(',', '.')

def fmt_pct(n):
    return f"{n:.1f}%"

# ============================================================
# 1. PPTX
# ============================================================
print("Gerando PPTX com identidade visual...")
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*hex_to_rgb(color))

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*hex_to_rgb(color))
    p.font.name = font_name
    p.alignment = align
    return tf

def add_logo(slide, left, top, width=1.2):
    """Adiciona logo ao slide"""
    try:
        slide.shapes.add_picture(LOGO_PATH, Inches(left), Inches(top), Inches(width), Inches(width))
    except Exception as e:
        print(f'  [WARN] Nao foi possivel adicionar logo: {e}')

def add_title_slide(title, subtitle=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, DARK_BG)
    # Logo centralizada no topo
    add_logo(slide, 5.8, 0.8, 1.6)
    add_text_box(slide, 1, 2.8, 11, 1.2, title, font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, 1, 4.0, 11, 1.2, subtitle, font_size=16, color=MUTED, align=PP_ALIGN.CENTER)
    # Linha decorativa
    shape = slide.shapes.add_shape(1, Inches(4), Inches(5.2), Inches(5.33), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb(ACCENT))
    shape.line.fill.background()
    return slide

def add_section_slide(title, subtitle=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, LIGHT_BG)
    # Logo pequena
    add_logo(slide, 0.5, 0.3, 0.7)
    add_text_box(slide, 0.5, 1.5, 12, 0.9, title, font_size=34, bold=True, color=PRIMARY)
    if subtitle:
        add_text_box(slide, 0.5, 2.3, 12, 0.5, subtitle, font_size=14, color=MUTED)
    # Linha
    shape = slide.shapes.add_shape(1, Inches(0.5), Inches(2.1), Inches(3), Inches(0.03))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb(ACCENT))
    shape.line.fill.background()
    return slide

def add_kpi_card(slide, left, top, width, height, title, value, subtitle='', color=PRIMARY):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb(WHITE))
    shape.line.fill.background()

    border = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(0.04))
    border.fill.solid()
    border.fill.fore_color.rgb = RGBColor(*hex_to_rgb(color))
    border.line.fill.background()

    add_text_box(slide, left + 0.2, top + 0.15, width - 0.4, 0.3, title, font_size=11, color=MUTED)
    add_text_box(slide, left + 0.2, top + 0.45, width - 0.4, 0.5, value, font_size=28, bold=True, color=color)
    if subtitle:
        add_text_box(slide, left + 0.2, top + 0.95, width - 0.4, 0.3, subtitle, font_size=9, color=MUTED)

# ===== SLIDE 1: CAPA =====
add_title_slide(
    'Relatório de Taxa de Ociosidade',
    'Análise de Viabilidade: Abertura 8h · Fechamento 19h · Sábados\nJaneiro a Junho 2026'
)

# ===== SLIDE 2: RESUMO EXECUTIVO =====
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide2, LIGHT_BG)
add_logo(slide2, 0.5, 0.3, 0.6)
add_text_box(slide2, 1.3, 0.35, 11, 0.6, 'Resumo Executivo', font_size=30, bold=True, color=PRIMARY)
add_text_box(slide2, 1.3, 0.85, 11, 0.3, 'Análise combinada dos 3 cenários de redução de horário — Janeiro a Junho 2026', font_size=12, color=MUTED)

add_kpi_card(slide2, 0.5, 1.5, 2.9, 1.1, 'TOTAL ATENDIMENTOS IMPACTADOS', str(kpis['total_geral']),
             'Soma de todos os cenários', PRIMARY)
add_kpi_card(slide2, 3.7, 1.5, 2.9, 1.1, 'SÁBADOS COM ATENDIMENTO', f"{kpis['sabados_unicos']}/{kpis['total_sabados_periodo']}",
             f"Taxa de ocupação: {fmt_pct(kpis['taxa_ocupacao_sabados'])}", ACCENT)
add_kpi_card(slide2, 7.0, 1.5, 2.9, 1.1, 'MÉDIA/DIA ANTES 8H', f"{kpis['media_antes_8h_dia']:.2f}",
             f'{kpis["total_antes_8h"]} atendimentos em {kpis["dias_uteis"]} dias úteis', SECONDARY)
add_kpi_card(slide2, 10.2, 1.5, 2.9, 1.1, 'ATENDIMENTOS APÓS 19H', str(kpis['total_apos_19h']),
             'Em 6 meses de análise', SUCCESS_COLOR)

add_text_box(slide2, 0.7, 3.0, 12, 0.4, 'Conclusão Preliminar', font_size=18, bold=True, color=PRIMARY)
conclusao = (
    'Abertura às 8h: VIÁVEL — apenas 0,51 atendimentos/dia antes das 8h (66 em 129 dias úteis). '
    'Impacto mínimo na operação.\n'
    'Fechamento às 19h: VIÁVEL — apenas 24 atendimentos após 19h em 6 meses. '
    'Maioria de estética (16 de 24).\n'
    'Sábados: REQUER ANÁLISE — 80,8% dos sábados têm atendimento, '
    'média de 12,4 atendimentos/dia, 260 no total.'
)
add_text_box(slide2, 0.9, 3.5, 11.5, 2.5, conclusao, font_size=13, color=TEXT)

# ===== SLIDE 3: SÁBADOS =====
slide3 = add_section_slide('Análise Detalhada: Sábados', '21 de 26 sábados com atendimento — 80,8% de ocupação')

add_kpi_card(slide3, 0.5, 2.8, 2.4, 1.0, 'TOTAL ATENDIMENTOS', str(kpis['total_sabados']), '', ACCENT)
add_kpi_card(slide3, 3.2, 2.8, 2.4, 1.0, 'MÉDIA/SÁBADO', f"{kpis['media_sabado_dia']:.1f}", '', ACCENT)
add_kpi_card(slide3, 5.9, 2.8, 2.4, 1.0, 'SÁBADOS OCUPADOS', f"{kpis['sabados_unicos']}/{kpis['total_sabados_periodo']}", '', ACCENT)
add_kpi_card(slide3, 8.6, 2.8, 2.4, 1.0, 'PACIENTES ÚNICOS', str(kpis['sabados_pacientes_unicos']), '', ACCENT)

chart_data = CategoryChartData()
chart_data.categories = list(kpis['sabados_por_mes'].keys())
chart_data.add_series('Atendimentos', list(kpis['sabados_por_mes'].values()))
cf = slide3.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(4.1), Inches(7.5), Inches(3.0), chart_data)
c = cf.chart
c.has_legend = False
plot = c.plots[0]
plot.gap_width = 80
series = plot.series[0]
series.format.fill.solid()
series.format.fill.fore_color.rgb = RGBColor(*hex_to_rgb(ACCENT))

add_text_box(slide3, 8.5, 4.1, 4.5, 0.3, 'Top 8 Tipos (Sábados)', font_size=14, bold=True, color=PRIMARY)
top_tipos = list(kpis['sabados_por_tipo'].items())[:8]
y = 4.6
for tipo, qtd in top_tipos:
    add_text_box(slide3, 8.7, y, 3.2, 0.28, tipo[:45], font_size=9, color=TEXT)
    add_text_box(slide3, 12.2, y, 0.6, 0.28, str(qtd), font_size=9, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
    y += 0.28

# ===== SLIDE 4: ANTES 8H =====
slide4 = add_section_slide('Análise Detalhada: Antes das 8h', 'Atendimentos entre 7h e 8h da manhã')

add_kpi_card(slide4, 0.5, 2.8, 2.4, 1.0, 'TOTAL', str(kpis['total_antes_8h']), '', SECONDARY)
add_kpi_card(slide4, 3.2, 2.8, 2.4, 1.0, 'MÉDIA/DIA ÚTIL', f"{kpis['media_antes_8h_dia']:.2f}", '', SECONDARY)
add_kpi_card(slide4, 5.9, 2.8, 2.4, 1.0, 'DIAS ÚTEIS', str(kpis['dias_uteis']), '', SECONDARY)
add_kpi_card(slide4, 8.6, 2.8, 2.4, 1.0, 'PACIENTES ÚNICOS', str(kpis['antes_8h_pacientes_unicos']), '', SECONDARY)

chart_data2 = CategoryChartData()
chart_data2.categories = list(kpis['antes_8h_por_mes'].keys())
chart_data2.add_series('Atendimentos', list(kpis['antes_8h_por_mes'].values()))
cf2 = slide4.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(4.1), Inches(7.5), Inches(3.0), chart_data2)
c2 = cf2.chart
c2.has_legend = False
plot2 = c2.plots[0]
plot2.gap_width = 80
series2 = plot2.series[0]
series2.format.fill.solid()
series2.format.fill.fore_color.rgb = RGBColor(*hex_to_rgb(SECONDARY))

add_text_box(slide4, 8.5, 4.1, 4.5, 0.3, 'Tipos de Atendimento (Antes 8h)', font_size=14, bold=True, color=PRIMARY)
y = 4.6
for tipo, qtd in kpis['antes_8h_por_tipo'].items():
    add_text_box(slide4, 8.7, y, 3.2, 0.28, tipo[:45], font_size=9, color=TEXT)
    add_text_box(slide4, 12.2, y, 0.6, 0.28, str(qtd), font_size=9, bold=True, color=SECONDARY, align=PP_ALIGN.RIGHT)
    y += 0.28

# ===== SLIDE 5: APÓS 19H =====
slide5 = add_section_slide('Análise Detalhada: Após 19h', 'Atendimentos após o horário comercial')

add_kpi_card(slide5, 0.5, 2.8, 3.0, 1.0, 'TOTAL', str(kpis['total_apos_19h']), '24 atendimentos em 6 meses', SUCCESS_COLOR)
add_kpi_card(slide5, 3.8, 2.8, 3.0, 1.0, 'PACIENTES ÚNICOS', str(kpis['apos_19h_pacientes_unicos']), '', SUCCESS_COLOR)
add_kpi_card(slide5, 7.1, 2.8, 3.0, 1.0, 'PRINCIPAL ÁREA', 'Estética (16/24)', '66,7% dos atendimentos', SUCCESS_COLOR)

chart_data3 = CategoryChartData()
chart_data3.categories = list(kpis['apos_19h_por_tipo'].keys())
chart_data3.add_series('Atendimentos', list(kpis['apos_19h_por_tipo'].values()))
cf3 = slide5.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.5), Inches(4.1), Inches(12.3), Inches(3.0), chart_data3)
c3 = cf3.chart
c3.has_legend = False
plot3 = c3.plots[0]
series3 = plot3.series[0]
series3.format.fill.solid()
series3.format.fill.fore_color.rgb = RGBColor(*hex_to_rgb(SUCCESS_COLOR))

# ===== SLIDE 6: IMPACTO POR PROFISSIONAL =====
slide6 = add_section_slide('Impacto por Profissional', 'Todos os cenários combinados')

add_text_box(slide6, 0.5, 2.8, 6, 0.4, 'Top Profissionais Impactados', font_size=16, bold=True, color=PRIMARY)
y = 3.4
max_val = max(kpis['top10_profissionais_impactados'].values())
for prof, qtd in kpis['top10_profissionais_impactados'].items():
    pct = qtd / kpis['total_geral'] * 100
    bar_w = (qtd / max_val) * 9
    bar = slide6.shapes.add_shape(1, Inches(2.5), Inches(y), Inches(bar_w), Inches(0.35))
    bar.fill.solid()
    bar_color = ACCENT if pct > 20 else SECONDARY
    bar.fill.fore_color.rgb = RGBColor(*hex_to_rgb(bar_color))
    bar.line.fill.background()
    add_text_box(slide6, 0.5, y, 1.9, 0.35, prof[:35], font_size=10, color=TEXT, align=PP_ALIGN.RIGHT)
    add_text_box(slide6, 2.5 + bar_w + 0.1, y, 1.5, 0.35, f'{qtd} ({fmt_pct(pct)})', font_size=10, bold=True, color=PRIMARY)
    y += 0.5

# ===== SLIDE 7: RECOMENDAÇÕES =====
slide7 = add_section_slide('Recomendações', 'Análise de viabilidade por cenário')

recs = [
    ('Abertura às 8h', 'VIÁVEL',
     'Apenas 66 atendimentos em 6 meses. Média de 0,51/dia útil. Implementar imediatamente.',
     SUCCESS_COLOR),
    ('Fechamento às 19h', 'VIÁVEL',
     'Apenas 24 atendimentos em 6 meses. Manter 1 dia/semana estendido para estética.',
     SUCCESS_COLOR),
    ('Eliminar Sábados', 'ATENÇÃO',
     '260 atendimentos, 80,8% ocupação. Reduzir para 2 sábados/mês. Reavaliar em 6 meses.',
     WARNING_COLOR),
    ('Cenário Combinado', 'PLANO NECESSÁRIO',
     '350 atendimentos impactados (~58/mês). Transição 30-60 dias. Meta: Setembro/2026.',
     DANGER_COLOR),
]

y = 2.8
for titulo, status, desc, cor in recs:
    badge = slide7.shapes.add_shape(1, Inches(0.5), Inches(y), Inches(2.0), Inches(0.35))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(*hex_to_rgb(cor))
    badge.line.fill.background()
    add_text_box(slide7, 0.5, y, 2.0, 0.35, status, font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide7, 2.7, y - 0.05, 10, 0.35, titulo, font_size=16, bold=True, color=PRIMARY)
    add_text_box(slide7, 2.7, y + 0.35, 10.3, 0.5, desc, font_size=11, color=TEXT)
    y += 0.85

# ===== SLIDE 8: CAPA FINAL =====
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide8, DARK_BG)
add_logo(slide8, 5.8, 0.8, 1.6)
add_text_box(slide8, 1, 2.5, 11, 1, 'Obrigado!', font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

passos = [
    '1. Validar análise com equipe',
    '2. Simular realocação dos atendimentos de sábado',
    '3. Definir cronograma de comunicação',
    '4. Período de transição: 30-60 dias',
    f'5. Meta: novo horário até Setembro/2026',
]
y = 3.8
for p in passos:
    add_text_box(slide8, 3, y, 7, 0.4, p, font_size=14, color='#d5c8e0', align=PP_ALIGN.LEFT)
    y += 0.45

shape = slide8.shapes.add_shape(1, Inches(4), Inches(6.5), Inches(5.33), Inches(0.03))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb(ACCENT))
shape.line.fill.background()

pptx_path = f'{BASE}/Relatorio_Taxa_Ociosidade_Relive.pptx'
prs.save(pptx_path)
print(f'[OK] PPTX salvo: {pptx_path}')

# ============================================================
# 2. PDF
# ============================================================
print("Gerando PDF com identidade visual...")
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import mm, cm
from reportlab.lib.colors import HexColor, white, black
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer, Table,
                                  TableStyle, PageBreak, Image as RLImage, KeepTogether)
from reportlab.graphics.shapes import Drawing, Rect, String, Line
from reportlab.graphics.charts.barcharts import VerticalBarChart, HorizontalBarChart

pdf_path = f'{BASE}/Relatorio_Taxa_Ociosidade_Relive.pdf'
doc = SimpleDocTemplate(pdf_path, pagesize=A4,
                        leftMargin=20*mm, rightMargin=20*mm,
                        topMargin=20*mm, bottomMargin=20*mm)

styles = getSampleStyleSheet()
styleTitle = ParagraphStyle('CustomTitle', parent=styles['Title'], fontSize=22,
                             textColor=HexColor(PRIMARY), spaceAfter=6*mm,
                             fontName='Helvetica-Bold')
styleH2 = ParagraphStyle('CustomH2', parent=styles['Heading2'], fontSize=16,
                          textColor=HexColor(PRIMARY), spaceAfter=4*mm, spaceBefore=8*mm,
                          fontName='Helvetica-Bold')
styleH3 = ParagraphStyle('CustomH3', parent=styles['Heading3'], fontSize=13,
                          textColor=HexColor(SECONDARY), spaceAfter=3*mm, spaceBefore=6*mm,
                          fontName='Helvetica-Bold')
styleBody = ParagraphStyle('CustomBody', parent=styles['Normal'], fontSize=10,
                            textColor=HexColor(TEXT), spaceAfter=3*mm, leading=15)
styleKpiLabel = ParagraphStyle('KPILabel', parent=styles['Normal'], fontSize=8,
                                textColor=HexColor(MUTED), alignment=TA_CENTER, spaceAfter=1*mm)
styleFooter = ParagraphStyle('Footer', parent=styles['Normal'], fontSize=8,
                              textColor=HexColor(MUTED), alignment=TA_CENTER)

elements = []

# CAPA
elements.append(Spacer(1, 25*mm))
try:
    logo_img = RLImage(LOGO_PATH, width=40*mm, height=40*mm)
    elements.append(logo_img)
except:
    pass
elements.append(Spacer(1, 10*mm))
elements.append(Paragraph('Relatório de Taxa de Ociosidade', styleTitle))
elements.append(Spacer(1, 5*mm))
elements.append(Paragraph('Análise de Viabilidade para Redução de Horário', styleH2))
elements.append(Paragraph('Abertura 8h · Fechamento 19h · Eliminação de Sábados', styleBody))
elements.append(Paragraph(f'Período: Janeiro a Junho 2026 | Relive', styleBody))
elements.append(Paragraph(f'Gerado em: {datetime.now().strftime("%d/%m/%Y %H:%M")}', styleBody))
elements.append(Spacer(1, 10*mm))

d = Drawing(170*mm, 2*mm)
d.add(Line(0, 0, 170*mm, 0, strokeColor=HexColor(ACCENT), strokeWidth=2))
elements.append(d)
elements.append(Spacer(1, 10*mm))

# RESUMO EXECUTIVO
elements.append(Paragraph('1. Resumo Executivo', styleH2))
elements.append(Paragraph(
    f'Este relatório analisa a viabilidade de três mudanças operacionais na clínica Relive: '
    f'<b>(a)</b> abertura às 8h em vez de 7h/7h30, <b>(b)</b> fechamento às 19h em vez de após 19h, '
    f'e <b>(c)</b> não abertura aos sábados. A análise compreende o período de 01/01/2026 a 30/06/2026.',
    styleBody
))

kpi_table_data = [
    ['Indicador', 'Valor', 'Interpretação'],
    ['Total atendimentos impactados', str(kpis['total_geral']), 'Soma dos 3 cenários (6 meses)'],
    ['Atendimentos antes 8h', str(kpis['total_antes_8h']), f'Média {kpis["media_antes_8h_dia"]:.2f}/dia útil'],
    ['Atendimentos após 19h', str(kpis['total_apos_19h']), 'Média 4/mês'],
    ['Atendimentos sábados', str(kpis['total_sabados']), f'{kpis["sabados_unicos"]}/{kpis["total_sabados_periodo"]} sábados ocupados'],
    ['Taxa ocupação sábados', fmt_pct(kpis['taxa_ocupacao_sabados']), f'Média {kpis["media_sabado_dia"]:.1f} atend/sábado'],
    ['Pacientes únicos impactados', str(kpis['sabados_pacientes_unicos'] + kpis['antes_8h_pacientes_unicos'] + kpis['apos_19h_pacientes_unicos']), 'Todos os cenários'],
]

kpi_table = Table(kpi_table_data, colWidths=[60*mm, 35*mm, 65*mm])
kpi_table.setStyle(TableStyle([
    ('BACKGROUND', (0, 0), (-1, 0), HexColor(PRIMARY)),
    ('TEXTCOLOR', (0, 0), (-1, 0), white),
    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
    ('FONTSIZE', (0, 0), (-1, 0), 10),
    ('FONTSIZE', (0, 1), (-1, -1), 9),
    ('ALIGN', (1, 0), (1, -1), 'CENTER'),
    ('GRID', (0, 0), (-1, -1), 0.5, HexColor('#d5c8e0')),
    ('ROWBACKGROUNDS', (0, 1), (-1, -1), [white, HexColor('#f8f4fc')]),
    ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
    ('TOPPADDING', (0, 0), (-1, -1), 4),
    ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
]))
elements.append(kpi_table)
elements.append(Spacer(1, 6*mm))

# Conclusão preliminar
elements.append(Paragraph('2. Conclusão Preliminar', styleH2))
for c in [
    '<b>Abertura às 8h: VIÁVEL</b> — Apenas 0,51 atendimentos/dia antes das 8h. Impacto mínimo.',
    '<b>Fechamento às 19h: VIÁVEL</b> — Apenas 24 atendimentos após 19h em 6 meses.',
    '<b>Sábados: REQUER ATENÇÃO</b> — 80,8% de ocupação, 260 atendimentos. Eliminar exige plano.',
]:
    elements.append(Paragraph(c, styleBody))

elements.append(PageBreak())

# SÁBADOS
elements.append(Paragraph('3. Análise Detalhada: Sábados', styleH2))

drawing = Drawing(170*mm, 80*mm)
bc = VerticalBarChart()
bc.x = 40; bc.y = 40; bc.height = 120; bc.width = 400
meses_sab = list(kpis['sabados_por_mes'].keys())
valores_sab = list(kpis['sabados_por_mes'].values())
bc.data = [valores_sab]
bc.categoryAxis.categoryNames = [m[-2:] for m in meses_sab]
bc.categoryAxis.labels.fontSize = 8
bc.valueAxis.valueMin = 0
bc.valueAxis.valueMax = max(valores_sab) * 1.2
bc.valueAxis.valueStep = 10
bc.bars[0].fillColor = HexColor(ACCENT)
bc.barWidth = 15; bc.barSpacing = 8
drawing.add(bc)
drawing.add(String(200, 175, 'Atendimentos aos Sábados por Mês (2026)', fontSize=11, fillColor=HexColor(PRIMARY), textAnchor='middle'))
elements.append(drawing)
elements.append(Spacer(1, 5*mm))

sab_table_data = [['Métrica', 'Valor']]
for m, v in [
    ('Total atendimentos', str(kpis['total_sabados'])),
    ('Sábados com atendimento', f'{kpis["sabados_unicos"]}/{kpis["total_sabados_periodo"]}'),
    ('Sábados sem atendimento', str(kpis['sabados_sem_atendimento'])),
    ('Taxa de ocupação', fmt_pct(kpis['taxa_ocupacao_sabados'])),
    ('Média atendimentos/sábado', f'{kpis["media_sabado_dia"]:.1f}'),
    ('Pacientes únicos', str(kpis['sabados_pacientes_unicos'])),
]:
    sab_table_data.append([m, v])

sab_table = Table(sab_table_data, colWidths=[80*mm, 60*mm])
sab_table.setStyle(TableStyle([
    ('BACKGROUND', (0, 0), (-1, 0), HexColor(ACCENT)),
    ('TEXTCOLOR', (0, 0), (-1, 0), white),
    ('GRID', (0, 0), (-1, -1), 0.5, HexColor('#d5c8e0')),
    ('ROWBACKGROUNDS', (0, 1), (-1, -1), [white, HexColor('#f8f4fc')]),
    ('FONTSIZE', (0, 0), (-1, -1), 9),
    ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
    ('TOPPADDING', (0, 0), (-1, -1), 3),
    ('BOTTOMPADDING', (0, 0), (-1, -1), 3),
]))
elements.append(sab_table)

elements.append(PageBreak())

# ANTES 8H
elements.append(Paragraph('4. Análise Detalhada: Antes das 8h', styleH2))

drawing2 = Drawing(170*mm, 80*mm)
bc2 = VerticalBarChart()
bc2.x = 40; bc2.y = 40; bc2.height = 120; bc2.width = 400
meses_8h = list(kpis['antes_8h_por_mes'].keys())
valores_8h = list(kpis['antes_8h_por_mes'].values())
bc2.data = [valores_8h]
bc2.categoryAxis.categoryNames = [m[-2:] for m in meses_8h]
bc2.categoryAxis.labels.fontSize = 8
bc2.valueAxis.valueMin = 0
bc2.valueAxis.valueMax = max(valores_8h) * 1.3 if valores_8h else 20
bc2.valueAxis.valueStep = 5
bc2.bars[0].fillColor = HexColor(SECONDARY)
bc2.barWidth = 15; bc2.barSpacing = 8
drawing2.add(bc2)
drawing2.add(String(200, 175, 'Atendimentos Antes 8h por Mês (2026)', fontSize=11, fillColor=HexColor(PRIMARY), textAnchor='middle'))
elements.append(drawing2)
elements.append(Spacer(1, 5*mm))

elements.append(Paragraph(
    f'Total de <b>{kpis["total_antes_8h"]}</b> atendimentos antes das 8h em {kpis["dias_uteis"]} dias úteis. '
    f'Média de <b>{kpis["media_antes_8h_dia"]:.2f}</b> atendimentos por dia útil.',
    styleBody
))

elements.append(PageBreak())

# APÓS 19H
elements.append(Paragraph('5. Análise Detalhada: Após 19h', styleH2))

drawing3 = Drawing(170*mm, 70*mm)
bc3 = HorizontalBarChart()
bc3.x = 50; bc3.y = 10; bc3.height = 100; bc3.width = 350
tipos_19h = [t[:40] for t in list(kpis['apos_19h_por_tipo'].keys())]
valores_19h = list(kpis['apos_19h_por_tipo'].values())
bc3.data = [valores_19h]
bc3.categoryAxis.categoryNames = tipos_19h
bc3.categoryAxis.labels.fontSize = 7
bc3.valueAxis.valueMin = 0
bc3.valueAxis.valueMax = max(valores_19h) + 2 if valores_19h else 8
bc3.valueAxis.valueStep = 1
bc3.bars[0].fillColor = HexColor(SUCCESS_COLOR)
drawing3.add(bc3)
drawing3.add(String(200, 120, 'Tipos de Atendimento Após 19h', fontSize=10, fillColor=HexColor(PRIMARY), textAnchor='middle'))
elements.append(drawing3)
elements.append(Spacer(1, 5*mm))

elements.append(Paragraph(
    f'Apenas <b>{kpis["total_apos_19h"]}</b> atendimentos após 19h em 6 meses. Maioria de estética.',
    styleBody
))

elements.append(PageBreak())

# IMPACTO POR PROFISSIONAL
elements.append(Paragraph('6. Impacto por Profissional', styleH2))

prof_table_data = [['Profissional', 'Atendimentos', '% do Total']]
for prof, qtd in kpis['top10_profissionais_impactados'].items():
    pct = qtd / kpis['total_geral'] * 100
    prof_table_data.append([prof, str(qtd), fmt_pct(pct)])

prof_table = Table(prof_table_data, colWidths=[80*mm, 40*mm, 40*mm])
prof_table.setStyle(TableStyle([
    ('BACKGROUND', (0, 0), (-1, 0), HexColor(PRIMARY)),
    ('TEXTCOLOR', (0, 0), (-1, 0), white),
    ('GRID', (0, 0), (-1, -1), 0.5, HexColor('#d5c8e0')),
    ('ROWBACKGROUNDS', (0, 1), (-1, -1), [white, HexColor('#f8f4fc')]),
    ('FONTSIZE', (0, 0), (-1, -1), 9),
    ('ALIGN', (1, 0), (-1, -1), 'CENTER'),
    ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
    ('TOPPADDING', (0, 0), (-1, -1), 4),
    ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
]))
elements.append(prof_table)
elements.append(Spacer(1, 8*mm))

# RECOMENDAÇÕES
elements.append(Paragraph('7. Recomendações', styleH2))
recs_text = [
    '<b>Abertura às 8h:</b> VIÁVEL. Implementar imediatamente. Rearranjar horários para 8h.',
    '<b>Fechamento às 19h:</b> VIÁVEL. Manter 1 dia/semana com horário estendido para estética.',
    '<b>Sábados:</b> NÃO ELIMINAR. Reduzir para 2 sábados/mês. Reavaliar após 6 meses.',
    '<b>Plano de transição:</b> 30-60 dias. Meta: novo horário até Setembro/2026.',
]
for r in recs_text:
    elements.append(Paragraph(r, styleBody))

elements.append(Spacer(1, 8*mm))
d2 = Drawing(170*mm, 2*mm)
d2.add(Line(0, 0, 170*mm, 0, strokeColor=HexColor(ACCENT), strokeWidth=1))
elements.append(d2)
elements.append(Paragraph('Relatório gerado automaticamente — Dados: Janeiro a Junho 2026 | Relive', styleFooter))

doc.build(elements)
print(f'[OK] PDF salvo: {pdf_path}')

# ============================================================
# 3. HTML
# ============================================================
print("Gerando HTML com identidade visual...")

# Converter logo para base64 para embed no HTML
import base64
logo_b64 = ''
try:
    with open(LOGO_PATH, 'rb') as lf:
        logo_b64 = base64.b64encode(lf.read()).decode('utf-8')
except:
    pass

html_content = f'''<!DOCTYPE html>
<html lang="pt-BR">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>Relatório de Taxa de Ociosidade — Relive</title>
<script src="https://cdn.jsdelivr.net/npm/chart.js@4.4.0/dist/chart.umd.min.js"></script>
<style>
* {{ margin: 0; padding: 0; box-sizing: border-box; }}
body {{ font-family: 'Segoe UI', system-ui, -apple-system, sans-serif; background: #f0eef5; color: {TEXT}; line-height: 1.6; }}
.header {{ background: linear-gradient(135deg, {DARK_BG} 0%, {PRIMARY} 50%, {SECONDARY} 100%); color: {WHITE}; padding: 50px 40px 40px; text-align: center; position: relative; overflow: hidden; }}
.header::before {{ content: ''; position: absolute; top: -50%; left: -50%; width: 200%; height: 200%; background: radial-gradient(circle, rgba(255,255,255,0.05) 1px, transparent 1px); background-size: 30px 30px; }}
.header .logo {{ width: 90px; height: 90px; margin-bottom: 15px; position: relative; z-index: 1; }}
.header h1 {{ font-size: 2.5em; margin-bottom: 8px; position: relative; z-index: 1; font-weight: 700; }}
.header p {{ font-size: 1.1em; opacity: 0.85; position: relative; z-index: 1; }}
.container {{ max-width: 1300px; margin: 0 auto; padding: 30px 20px; }}
.kpi-grid {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(280px, 1fr)); gap: 20px; margin-bottom: 40px; }}
.kpi-card {{ background: {WHITE}; border-radius: 12px; padding: 25px; box-shadow: 0 2px 12px rgba(96,32,96,0.08); border-top: 4px solid {PRIMARY}; transition: transform 0.2s; }}
.kpi-card:hover {{ transform: translateY(-2px); box-shadow: 0 4px 20px rgba(96,32,96,0.15); }}
.kpi-card.warning {{ border-top-color: {ACCENT}; }}
.kpi-card.success {{ border-top-color: {SUCCESS_COLOR}; }}
.kpi-card.info {{ border-top-color: {SECONDARY}; }}
.kpi-card .kpi-label {{ font-size: 0.8em; color: {MUTED}; text-transform: uppercase; letter-spacing: 1px; margin-bottom: 8px; font-weight: 600; }}
.kpi-card .kpi-value {{ font-size: 2.2em; font-weight: 700; color: {PRIMARY}; }}
.kpi-card.warning .kpi-value {{ color: {ACCENT}; }}
.kpi-card.success .kpi-value {{ color: {SUCCESS_COLOR}; }}
.kpi-card.info .kpi-value {{ color: {SECONDARY}; }}
.kpi-card .kpi-sub {{ font-size: 0.85em; color: {MUTED}; margin-top: 5px; }}
.section {{ background: {WHITE}; border-radius: 12px; padding: 30px; margin-bottom: 30px; box-shadow: 0 2px 12px rgba(96,32,96,0.08); }}
.section h2 {{ font-size: 1.5em; color: {PRIMARY}; margin-bottom: 15px; padding-bottom: 10px; border-bottom: 2px solid #e8dff0; display: flex; align-items: center; gap: 10px; }}
.section h3 {{ font-size: 1.15em; color: {SECONDARY}; margin: 20px 0 10px; }}
.chart-row {{ display: grid; grid-template-columns: 1fr 1fr; gap: 30px; margin: 20px 0; }}
.chart-box {{ background: #fdfbff; border-radius: 8px; padding: 20px; border: 1px solid #e8dff0; }}
.chart-box.full {{ grid-column: 1 / -1; }}
.chart-box canvas {{ max-height: 350px; }}
.conclusion-box {{ padding: 20px 25px; border-radius: 0 10px 10px 0; margin: 15px 0; border-left: 5px solid {PRIMARY}; }}
.conclusion-box.warning {{ background: #fff8f0; border-left-color: {ACCENT}; }}
.conclusion-box.success {{ background: #f0faf2; border-left-color: {SUCCESS_COLOR}; }}
.conclusion-box h4 {{ margin-bottom: 6px; font-size: 1.05em; }}
table {{ width: 100%; border-collapse: collapse; margin: 15px 0; }}
th {{ background: {PRIMARY}; color: {WHITE}; padding: 12px 15px; text-align: left; font-size: 0.85em; font-weight: 600; text-transform: uppercase; letter-spacing: 0.5px; }}
td {{ padding: 10px 15px; border-bottom: 1px solid #e8dff0; font-size: 0.9em; }}
tr:nth-child(even) {{ background: #fdfbff; }}
tr:hover {{ background: #f5f0f8; }}
.badge {{ display: inline-block; padding: 4px 14px; border-radius: 20px; font-size: 0.75em; font-weight: 700; text-transform: uppercase; letter-spacing: 0.5px; }}
.badge-ok {{ background: #d5f5e3; color: #1e7e34; }}
.badge-warn {{ background: #fef0d0; color: #b8730e; }}
.badge-alert {{ background: #fadbd8; color: #c0392b; }}
.footer {{ text-align: center; padding: 30px; color: {MUTED}; font-size: 0.85em; background: {DARK_BG}; margin-top: 40px; }}
.footer p {{ color: #b8a8c8; }}
.progress-bar {{ height: 8px; background: #e8dff0; border-radius: 4px; margin-top: 5px; overflow: hidden; }}
.progress-fill {{ height: 100%; border-radius: 4px; transition: width 1s ease; }}
.logo-small {{ width: 40px; height: 40px; vertical-align: middle; }}
@media (max-width: 768px) {{ .chart-row {{ grid-template-columns: 1fr; }} .kpi-grid {{ grid-template-columns: 1fr; }} .header h1 {{ font-size: 1.6em; }} }}
</style>
</head>
<body>

<div class="header">
    <img src="data:image/png;base64,{logo_b64}" alt="Relive" class="logo">
    <h1>Relatório de Taxa de Ociosidade</h1>
    <p>Análise de Viabilidade: Abertura 8h · Fechamento 19h · Sábados</p>
    <p style="margin-top:8px;font-size:0.9em;opacity:0.7;">Período: Janeiro a Junho 2026</p>
</div>

<div class="container">

    <div class="kpi-grid">
        <div class="kpi-card warning">
            <div class="kpi-label">Total Atendimentos Impactados</div>
            <div class="kpi-value">{kpis['total_geral']}</div>
            <div class="kpi-sub">Soma dos 3 cenários (6 meses)</div>
        </div>
        <div class="kpi-card">
            <div class="kpi-label">Sábados com Atendimento</div>
            <div class="kpi-value">{kpis['sabados_unicos']}/{kpis['total_sabados_periodo']}</div>
            <div class="kpi-sub">Taxa de ocupação: {fmt_pct(kpis['taxa_ocupacao_sabados'])}</div>
            <div class="progress-bar"><div class="progress-fill" style="width:{kpis['taxa_ocupacao_sabados']}%;background:{ACCENT};"></div></div>
        </div>
        <div class="kpi-card info">
            <div class="kpi-label">Média/Dia Antes 8h</div>
            <div class="kpi-value">{kpis['media_antes_8h_dia']:.2f}</div>
            <div class="kpi-sub">{kpis['total_antes_8h']} atendimentos em {kpis['dias_uteis']} dias úteis</div>
        </div>
        <div class="kpi-card success">
            <div class="kpi-label">Atendimentos Após 19h</div>
            <div class="kpi-value">{kpis['total_apos_19h']}</div>
            <div class="kpi-sub">Em 6 meses de análise</div>
        </div>
    </div>

    <div class="section">
        <h2>Conclusão Preliminar</h2>
        <div class="conclusion-box success">
            <h4>Abertura às 8h: VIÁVEL</h4>
            <p>Apenas 0,51 atendimentos/dia antes das 8h (66 em 129 dias úteis). Impacto mínimo. Profissionais: Helaine (42) e Cibele (22).</p>
        </div>
        <div class="conclusion-box success">
            <h4>Fechamento às 19h: VIÁVEL</h4>
            <p>Apenas 24 atendimentos após 19h em 6 meses. Maioria de estética (16 de 24), possível reacomodar.</p>
        </div>
        <div class="conclusion-box warning">
            <h4>Sábados: REQUER ANÁLISE CUIDADOSA</h4>
            <p>80,8% dos sábados têm atendimento, média 12,4/dia, 260 no total. Eliminar = redistribuir ~43 atendimentos/mês.</p>
        </div>
        <div class="conclusion-box" style="background:#fef5f5;border-left-color:{DANGER_COLOR};">
            <h4>Cenário Combinado: REQUER PLANO</h4>
            <p>350 atendimentos impactados em 6 meses (~58/mês). Transição 30-60 dias. Meta: Setembro/2026.</p>
        </div>
    </div>

    <div class="section">
        <h2>Sábados</h2>
        <div class="kpi-grid" style="margin-bottom:20px;">
            <div class="kpi-card warning">
                <div class="kpi-label">Total Atendimentos</div>
                <div class="kpi-value">{kpis['total_sabados']}</div>
            </div>
            <div class="kpi-card warning">
                <div class="kpi-label">Média por Sábado</div>
                <div class="kpi-value">{kpis['media_sabado_dia']:.1f}</div>
            </div>
            <div class="kpi-card warning">
                <div class="kpi-label">Sábados sem Atendimento</div>
                <div class="kpi-value">{kpis['sabados_sem_atendimento']}/{kpis['total_sabados_periodo']}</div>
            </div>
            <div class="kpi-card warning">
                <div class="kpi-label">Pacientes Únicos</div>
                <div class="kpi-value">{kpis['sabados_pacientes_unicos']}</div>
            </div>
        </div>
        <div class="chart-row">
            <div class="chart-box">
                <h3>Atendimentos por Mês</h3>
                <canvas id="chartSabadosMes"></canvas>
            </div>
            <div class="chart-box">
                <h3>Top 8 Tipos de Atendimento</h3>
                <canvas id="chartSabadosTipo"></canvas>
            </div>
        </div>
    </div>

    <div class="section">
        <h2>Antes das 8h</h2>
        <div class="kpi-grid" style="margin-bottom:20px;">
            <div class="kpi-card info">
                <div class="kpi-label">Total</div>
                <div class="kpi-value">{kpis['total_antes_8h']}</div>
            </div>
            <div class="kpi-card info">
                <div class="kpi-label">Média por Dia Útil</div>
                <div class="kpi-value">{kpis['media_antes_8h_dia']:.2f}</div>
            </div>
            <div class="kpi-card info">
                <div class="kpi-label">Dias Úteis no Período</div>
                <div class="kpi-value">{kpis['dias_uteis']}</div>
            </div>
            <div class="kpi-card info">
                <div class="kpi-label">Pacientes Únicos</div>
                <div class="kpi-value">{kpis['antes_8h_pacientes_unicos']}</div>
            </div>
        </div>
        <div class="chart-box full">
            <h3>Atendimentos por Mês</h3>
            <canvas id="chartAntes8hMes"></canvas>
        </div>
    </div>

    <div class="section">
        <h2>Após 19h</h2>
        <div class="kpi-grid" style="margin-bottom:20px;">
            <div class="kpi-card success">
                <div class="kpi-label">Total</div>
                <div class="kpi-value">{kpis['total_apos_19h']}</div>
            </div>
            <div class="kpi-card success">
                <div class="kpi-label">Pacientes Únicos</div>
                <div class="kpi-value">{kpis['apos_19h_pacientes_unicos']}</div>
            </div>
            <div class="kpi-card success">
                <div class="kpi-label">Principal Área</div>
                <div class="kpi-value" style="font-size:1.4em;">Estética</div>
                <div class="kpi-sub">16/24 atendimentos (66,7%)</div>
            </div>
        </div>
        <div class="chart-box">
            <h3>Tipos de Atendimento</h3>
            <canvas id="chartApos19hTipo"></canvas>
        </div>
    </div>

    <div class="section">
        <h2>Impacto por Profissional</h2>
        <div class="chart-box full" style="margin-bottom:20px;">
            <canvas id="chartProfissionais"></canvas>
        </div>
        <table>
            <tr><th>Profissional</th><th>Atendimentos</th><th>% do Total</th><th>Avaliação</th></tr>
'''

for prof, qtd in kpis['top10_profissionais_impactados'].items():
    pct = qtd / kpis['total_geral'] * 100
    badge_class = 'badge-alert' if pct > 30 else ('badge-warn' if pct > 10 else 'badge-ok')
    label = "ALTO" if pct > 30 else ("MEDIO" if pct > 10 else "BAIXO")
    html_content += f'<tr><td>{prof}</td><td><b>{qtd}</b></td><td>{fmt_pct(pct)}</td><td><span class="badge {badge_class}">{label}</span></td></tr>\n'

html_content += f'''
        </table>
    </div>

    <div class="section">
        <h2>Distribuição por Horário (Sábados)</h2>
        <div class="chart-box full">
            <canvas id="chartHorasSabado"></canvas>
        </div>
    </div>

    <div class="section">
        <h2>Recomendações</h2>
        <table>
            <tr><th>Cenário</th><th>Viabilidade</th><th>Recomendação</th></tr>
            <tr><td><b>Abertura às 8h</b></td><td><span class="badge badge-ok">VIÁVEL</span></td><td>Implementar imediatamente. Rearranjar horários para 8h.</td></tr>
            <tr><td><b>Fechamento às 19h</b></td><td><span class="badge badge-ok">VIÁVEL</span></td><td>Manter 1 dia/semana estendido para estética.</td></tr>
            <tr><td><b>Eliminar Sábados</b></td><td><span class="badge badge-warn">ATENÇÃO</span></td><td>Reduzir para 2 sábados/mês. Reavaliar em 6 meses.</td></tr>
            <tr><td><b>Cenário Combinado</b></td><td><span class="badge badge-alert">PLANO NECESSÁRIO</span></td><td>350 atendimentos impactados. Transição 30-60 dias. Meta: Setembro/2026.</td></tr>
        </table>
    </div>
</div>

<div class="footer">
    <p>Relatório de Taxa de Ociosidade — Relive</p>
    <p>Gerado em {datetime.now().strftime("%d/%m/%Y %H:%M")} | Dados: Janeiro a Junho 2026</p>
</div>

<script>
const primary = '{PRIMARY}';
const secondary = '{SECONDARY}';
const accent = '{ACCENT}';
const successColor = '{SUCCESS_COLOR}';

Chart.defaults.font.family = "'Segoe UI', system-ui, sans-serif";
Chart.defaults.font.size = 12;

new Chart(document.getElementById('chartSabadosMes'), {{
    type: 'bar',
    data: {{
        labels: {json.dumps([m[-2:] + '/' + m[:4] for m in kpis['sabados_por_mes'].keys()])},
        datasets: [{{ label: 'Atendimentos', data: {json.dumps(list(kpis['sabados_por_mes'].values()))}, backgroundColor: accent, borderRadius: 6 }}]
    }},
    options: {{ responsive: true, plugins: {{ legend: {{ display: false }} }}, scales: {{ y: {{ beginAtZero: true, ticks: {{ stepSize: 10 }} }} }} }}
}});

new Chart(document.getElementById('chartSabadosTipo'), {{
    type: 'bar',
    data: {{
        labels: {json.dumps([t[:30] for t in list(kpis['sabados_por_tipo'].keys())[:8]])},
        datasets: [{{ label: 'Atendimentos', data: {json.dumps(list(kpis['sabados_por_tipo'].values())[:8])}, backgroundColor: accent + 'cc', borderRadius: 6 }}]
    }},
    options: {{ indexAxis: 'y', responsive: true, plugins: {{ legend: {{ display: false }} }}, scales: {{ x: {{ beginAtZero: true, ticks: {{ stepSize: 5 }} }} }} }}
}});

new Chart(document.getElementById('chartAntes8hMes'), {{
    type: 'bar',
    data: {{
        labels: {json.dumps([m[-2:] + '/' + m[:4] for m in kpis['antes_8h_por_mes'].keys()])},
        datasets: [{{ label: 'Atendimentos', data: {json.dumps(list(kpis['antes_8h_por_mes'].values()))}, backgroundColor: secondary, borderRadius: 6 }}]
    }},
    options: {{ responsive: true, plugins: {{ legend: {{ display: false }} }}, scales: {{ y: {{ beginAtZero: true, ticks: {{ stepSize: 5 }} }} }} }}
}});

new Chart(document.getElementById('chartApos19hTipo'), {{
    type: 'bar',
    data: {{
        labels: {json.dumps([t[:30] for t in kpis['apos_19h_por_tipo'].keys()])},
        datasets: [{{ label: 'Atendimentos', data: {json.dumps(list(kpis['apos_19h_por_tipo'].values()))}, backgroundColor: successColor, borderRadius: 6 }}]
    }},
    options: {{ indexAxis: 'y', responsive: true, plugins: {{ legend: {{ display: false }} }}, scales: {{ x: {{ beginAtZero: true, ticks: {{ stepSize: 1 }} }} }} }}
}});

new Chart(document.getElementById('chartProfissionais'), {{
    type: 'bar',
    data: {{
        labels: {json.dumps([p[:30] for p in kpis['top10_profissionais_impactados'].keys()])},
        datasets: [{{ label: 'Atendimentos Impactados', data: {json.dumps(list(kpis['top10_profissionais_impactados'].values()))}, backgroundColor: [accent, accent, secondary, secondary, successColor, successColor, primary, primary], borderRadius: 6 }}]
    }},
    options: {{ responsive: true, plugins: {{ legend: {{ display: false }} }}, scales: {{ y: {{ beginAtZero: true }} }} }}
}});

new Chart(document.getElementById('chartHorasSabado'), {{
    type: 'line',
    data: {{
        labels: {json.dumps(list(kpis['sabados_por_hora'].keys()))},
        datasets: [{{ label: 'Atendimentos por Horário', data: {json.dumps(list(kpis['sabados_por_hora'].values()))}, borderColor: accent, backgroundColor: accent + '33', fill: true, tension: 0.3, pointRadius: 5, pointBackgroundColor: accent }}]
    }},
    options: {{ responsive: true, plugins: {{ legend: {{ display: false }} }}, scales: {{ y: {{ beginAtZero: true }} }} }}
}});
</script>
</body>
</html>
'''

html_path = f'{BASE}/Relatorio_Taxa_Ociosidade_Relive.html'
with open(html_path, 'w', encoding='utf-8') as f:
    f.write(html_content)
print(f'[OK] HTML salvo: {html_path}')

print('\n=== TODOS OS RELATORIOS REGERADOS COM IDENTIDADE VISUAL ===')
print(f'[PPTX] {pptx_path}')
print(f'[PDF]  {pdf_path}')
print(f'[HTML] {html_path}')
print(f'\nPaleta: Primary={PRIMARY} | Secondary={SECONDARY} | Accent={ACCENT}')
