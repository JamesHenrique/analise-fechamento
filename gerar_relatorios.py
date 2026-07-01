"""
Gerar Relatórios: PPTX, PDF, HTML
Relatório de Taxa de Ociosidade - Relive
"""
import json
from datetime import datetime

BASE = 'G:/Outros computadores/NOTEBOOK1/James/projeto_relive/Relatório/Relatório de abertura-fechamento'

with open(f'{BASE}/kpis.json', 'r', encoding='utf-8') as f:
    kpis = json.load(f)

# ============================================================
# CORES E ESTILOS PADRÃO
# ============================================================
COR_PRIMARY = '#1a5276'    # Azul escuro
COR_SECONDARY = '#2e86c1'  # Azul médio
COR_ACCENT = '#e74c3c'     # Vermelho destaque
COR_GREEN = '#27ae60'      # Verde
COR_ORANGE = '#f39c12'     # Laranja
COR_BG = '#f8f9fa'         # Background claro
COR_TEXT = '#2c3e50'       # Texto escuro

def fmt_num(n):
    return f"{n:,}".replace(',', '.')

def fmt_pct(n):
    return f"{n:.1f}%"

# ============================================================
# 1. GERAR PPTX
# ============================================================
print("Gerando PPTX...")
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 widescreen
prs.slide_height = Inches(7.5)

def add_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*hex_to_rgb(color))

def hex_to_rgb(h):
    h = h.lstrip('#')
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=COR_TEXT, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*hex_to_rgb(color))
    p.font.name = font_name
    p.alignment = align
    return tf

def add_title_slide(title, subtitle=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_slide_bg(slide, COR_PRIMARY)
    add_text_box(slide, 1, 2, 11, 1.5, title, font_size=40, bold=True, color='#ffffff', align=PP_ALIGN.CENTER)
    if subtitle:
        add_text_box(slide, 1, 3.5, 11, 1, subtitle, font_size=20, color='#bdc3c7', align=PP_ALIGN.CENTER)
    # Linha decorativa
    shape = slide.shapes.add_shape(1, Inches(4), Inches(3.2), Inches(5.33), Inches(0.03))  # retângulo
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb(COR_ACCENT))
    shape.line.fill.background()
    return slide

def add_kpi_card(slide, left, top, width, height, title, value, subtitle='', color=COR_PRIMARY):
    # Card background
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb('#ffffff'))
    shape.line.fill.background()
    shape.shadow.inherit = False

    # Borda superior colorida
    border = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(0.04))
    border.fill.solid()
    border.fill.fore_color.rgb = RGBColor(*hex_to_rgb(color))
    border.line.fill.background()

    add_text_box(slide, left + 0.2, top + 0.15, width - 0.4, 0.3, title, font_size=11, color='#7f8c8d')
    add_text_box(slide, left + 0.2, top + 0.45, width - 0.4, 0.5, value, font_size=28, bold=True, color=color)
    if subtitle:
        add_text_box(slide, left + 0.2, top + 0.95, width - 0.4, 0.3, subtitle, font_size=9, color='#95a5a6')

# ===== SLIDE 1: CAPA =====
add_title_slide(
    'Relatório de Taxa de Ociosidade',
    'Análise de Viabilidade: Abertura 8h · Fechamento 19h · Sábados\nPeríodo: Janeiro a Junho 2026 | Relive'
)

# ===== SLIDE 2: RESUMO EXECUTIVO =====
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide2, COR_BG)
add_text_box(slide2, 0.5, 0.3, 12, 0.7, 'Resumo Executivo', font_size=32, bold=True, color=COR_PRIMARY)
add_text_box(slide2, 0.5, 0.9, 12, 0.4, 'Análise combinada dos 3 cenários de redução de horário — Janeiro a Junho 2026', font_size=14, color='#7f8c8d')

# KPIs principais
add_kpi_card(slide2, 0.5, 1.6, 2.9, 1.2, 'TOTAL ATENDIMENTOS IMPACTADOS', str(kpis['total_geral']),
             'Soma de todos os cenários', COR_PRIMARY)
add_kpi_card(slide2, 3.7, 1.6, 2.9, 1.2, 'SÁBADOS COM ATENDIMENTO', f"{kpis['sabados_unicos']}/{kpis['total_sabados_periodo']}",
             f"Taxa de ocupação: {fmt_pct(kpis['taxa_ocupacao_sabados'])}", COR_ACCENT)
add_kpi_card(slide2, 7.0, 1.6, 2.9, 1.2, 'MÉDIA/DIA ANTES 8H', f"{kpis['media_antes_8h_dia']:.2f}",
             f'{kpis["total_antes_8h"]} atendimentos em {kpis["dias_uteis"]} dias úteis', COR_ORANGE)
add_kpi_card(slide2, 10.2, 1.6, 2.9, 1.2, 'ATENDIMENTOS APÓS 19H', str(kpis['total_apos_19h']),
             'Em 6 meses de análise', COR_GREEN)

# Recomendação
add_text_box(slide2, 0.5, 3.2, 12, 0.5, 'Conclusão Preliminar', font_size=20, bold=True, color=COR_PRIMARY)
conclusao = (
    '✓ Abertura às 8h: VIÁVEL — apenas 0,51 atendimentos/dia antes das 8h (66 em 129 dias úteis). '
    'Impacto mínimo na operação, concentrado em nutrição (Helaine e Cibele).\n'
    '✓ Fechamento às 19h: VIÁVEL — apenas 24 atendimentos após 19h em 6 meses. '
    'Maioria de estética (16 de 24), possível reacomodar.\n'
    '⚠ Sábados: REQUER ANÁLISE CUIDADOSA — 80,8% dos sábados têm atendimento, '
    'média de 12,4 atendimentos/dia, 260 no total. Eliminar exigiria redistribuir ~43 atendimentos/mês.'
)
add_text_box(slide2, 0.7, 3.8, 11.8, 2.8, conclusao, font_size=13, color=COR_TEXT)

# ===== SLIDE 3: SÁBADOS - ANÁLISE DETALHADA =====
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide3, COR_BG)
add_text_box(slide3, 0.5, 0.3, 12, 0.7, 'Análise Detalhada: Sábados', font_size=32, bold=True, color=COR_ACCENT)

# Cards
add_kpi_card(slide3, 0.5, 1.2, 2.4, 1.0, 'TOTAL ATENDIMENTOS', str(kpis['total_sabados']), '', COR_ACCENT)
add_kpi_card(slide3, 3.2, 1.2, 2.4, 1.0, 'MÉDIA/SÁBADO', f"{kpis['media_sabado_dia']:.1f}", '', COR_ACCENT)
add_kpi_card(slide3, 5.9, 1.2, 2.4, 1.0, 'SÁBADOS OCUPADOS', f"{kpis['sabados_unicos']}/{kpis['total_sabados_periodo']}", '', COR_ACCENT)
add_kpi_card(slide3, 8.6, 1.2, 2.4, 1.0, 'PACIENTES ÚNICOS', str(kpis['sabados_pacientes_unicos']), '', COR_ACCENT)

# Gráfico: atendimentos por mês (sábados)
chart_data = CategoryChartData()
chart_data.categories = list(kpis['sabados_por_mes'].keys())
chart_data.add_series('Atendimentos', list(kpis['sabados_por_mes'].values()))
chart_frame = slide3.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(2.5), Inches(6), Inches(4.2), chart_data)
chart = chart_frame.chart
chart.has_legend = False
plot = chart.plots[0]
plot.gap_width = 80
series = plot.series[0]
series.format.fill.solid()
series.format.fill.fore_color.rgb = RGBColor(*hex_to_rgb(COR_ACCENT))

# Top tipos (sábados)
add_text_box(slide3, 7, 2.5, 5.5, 0.4, 'Top 8 Tipos de Atendimento (Sábados)', font_size=14, bold=True, color=COR_PRIMARY)
top_tipos_sab = list(kpis['sabados_por_tipo'].items())[:8]
y = 3.0
for tipo, qtd in top_tipos_sab:
    add_text_box(slide3, 7.2, y, 4, 0.3, tipo[:50], font_size=10, color=COR_TEXT)
    add_text_box(slide3, 11.5, y, 1, 0.3, str(qtd), font_size=10, bold=True, color=COR_ACCENT, align=PP_ALIGN.RIGHT)
    y += 0.32

# ===== SLIDE 4: ANTES 8H =====
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide4, COR_BG)
add_text_box(slide4, 0.5, 0.3, 12, 0.7, 'Análise Detalhada: Atendimentos Antes das 8h', font_size=32, bold=True, color=COR_ORANGE)

add_kpi_card(slide4, 0.5, 1.2, 2.4, 1.0, 'TOTAL', str(kpis['total_antes_8h']), '', COR_ORANGE)
add_kpi_card(slide4, 3.2, 1.2, 2.4, 1.0, 'MÉDIA/DIA ÚTIL', f"{kpis['media_antes_8h_dia']:.2f}", '', COR_ORANGE)
add_kpi_card(slide4, 5.9, 1.2, 2.4, 1.0, 'DIAS ÚTEIS', str(kpis['dias_uteis']), '', COR_ORANGE)
add_kpi_card(slide4, 8.6, 1.2, 2.4, 1.0, 'PACIENTES ÚNICOS', str(kpis['antes_8h_pacientes_unicos']), '', COR_ORANGE)

# Gráfico mensal
chart_data2 = CategoryChartData()
chart_data2.categories = list(kpis['antes_8h_por_mes'].keys())
chart_data2.add_series('Atendimentos', list(kpis['antes_8h_por_mes'].values()))
chart_frame2 = slide4.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(2.5), Inches(6), Inches(4.2), chart_data2)
chart2 = chart_frame2.chart
chart2.has_legend = False
plot2 = chart2.plots[0]
plot2.gap_width = 80
series2 = plot2.series[0]
series2.format.fill.solid()
series2.format.fill.fore_color.rgb = RGBColor(*hex_to_rgb(COR_ORANGE))

# Top tipos antes 8h
add_text_box(slide4, 7, 2.5, 5.5, 0.4, 'Tipos de Atendimento (Antes 8h)', font_size=14, bold=True, color=COR_PRIMARY)
top_tipos_8h = list(kpis['antes_8h_por_tipo'].items())
y = 3.0
for tipo, qtd in top_tipos_8h:
    add_text_box(slide4, 7.2, y, 4, 0.3, tipo[:50], font_size=10, color=COR_TEXT)
    add_text_box(slide4, 11.5, y, 1, 0.3, str(qtd), font_size=10, bold=True, color=COR_ORANGE, align=PP_ALIGN.RIGHT)
    y += 0.32

# ===== SLIDE 5: APÓS 19H =====
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide5, COR_BG)
add_text_box(slide5, 0.5, 0.3, 12, 0.7, 'Análise Detalhada: Atendimentos Após 19h', font_size=32, bold=True, color=COR_GREEN)

add_kpi_card(slide5, 0.5, 1.2, 3.0, 1.0, 'TOTAL', str(kpis['total_apos_19h']), '24 atendimentos em 6 meses', COR_GREEN)
add_kpi_card(slide5, 3.8, 1.2, 3.0, 1.0, 'PACIENTES ÚNICOS', str(kpis['apos_19h_pacientes_unicos']), '', COR_GREEN)
add_kpi_card(slide5, 7.1, 1.2, 3.0, 1.0, 'PRINCIPAL ÁREA', 'Estética (16/24)', '66,7% dos atendimentos', COR_GREEN)

# Gráfico
chart_data3 = CategoryChartData()
chart_data3.categories = list(kpis['apos_19h_por_tipo'].keys())
chart_data3.add_series('Atendimentos', list(kpis['apos_19h_por_tipo'].values()))
chart_frame3 = slide5.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.5), Inches(2.5), Inches(12), Inches(4.5), chart_data3)
chart3 = chart_frame3.chart
chart3.has_legend = False
plot3 = chart3.plots[0]
series3 = plot3.series[0]
series3.format.fill.solid()
series3.format.fill.fore_color.rgb = RGBColor(*hex_to_rgb(COR_GREEN))

# ===== SLIDE 6: IMPACTO POR PROFISSIONAL =====
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide6, COR_BG)
add_text_box(slide6, 0.5, 0.3, 12, 0.7, 'Impacto por Profissional', font_size=32, bold=True, color=COR_PRIMARY)

add_text_box(slide6, 0.5, 1.2, 6, 0.4, 'Top Profissionais Impactados (Todos os Cenários)', font_size=14, bold=True, color=COR_PRIMARY)
y = 1.7
for prof, qtd in kpis['top10_profissionais_impactados'].items():
    pct = qtd / kpis['total_geral'] * 100
    # Barra
    bar_w = (qtd / max(kpis['top10_profissionais_impactados'].values())) * 8
    bar = slide6.shapes.add_shape(1, Inches(2.5), Inches(y), Inches(bar_w), Inches(0.35))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(*hex_to_rgb(COR_SECONDARY))
    bar.line.fill.background()
    add_text_box(slide6, 0.5, y, 2, 0.35, prof[:35], font_size=10, color=COR_TEXT, align=PP_ALIGN.RIGHT)
    add_text_box(slide6, 2.5 + bar_w + 0.1, y, 1.5, 0.35, f'{qtd} ({fmt_pct(pct)})', font_size=10, bold=True, color=COR_PRIMARY)
    y += 0.5

# ===== SLIDE 7: RECOMENDAÇÕES =====
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide7, COR_BG)
add_text_box(slide7, 0.5, 0.3, 12, 0.7, 'Recomendações', font_size=32, bold=True, color=COR_PRIMARY)

recomendacoes = [
    ('Abertura às 8h (eliminar 7h-7h30)', 'VIÁVEL',
     'Apenas 66 atendimentos em 6 meses. Média de 0,51/dia útil. Baixíssimo impacto. '
     'Profissionais afetados: Helaine (42) e Cibele (22). Rearranjar primeiros horários para 8h.',
     COR_GREEN),
    ('Fechamento às 19h (eliminar após 19h)', 'VIÁVEL',
     'Apenas 24 atendimentos em 6 meses. Maioria estética (16). Reacomodar dentro do horário comercial '
     'ou manter apenas dias específicos para estética noturna.',
     COR_GREEN),
    ('Eliminar abertura aos sábados', 'ATENÇÃO',
     '80,8% dos sábados têm atendimento, 260 no total, média 12,4/dia. Impacto ALTO. '
     'Alternativa: reduzir sábados para 2x/mês (manter apenas 1º e 3º sábado), '
     'reduzindo ~50% dos atendimentos de sábado sem eliminar completamente.',
     COR_ORANGE),
    ('Cenário combinado (8h-19h, sem sábados)', 'REQUER PLANO',
     '350 atendimentos impactados em 6 meses (~58/mês). Necessário plano de transição: '
     '1) Migrar sábados para dias úteis; 2) Reacomodar estética noturna; '
     '3) Comunicar pacientes com 30 dias de antecedência.',
     COR_ACCENT),
]

y = 1.3
for titulo, status, desc, cor_hex in recomendacoes:
    # Status badge
    badge = slide7.shapes.add_shape(1, Inches(0.5), Inches(y), Inches(1.6), Inches(0.35))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(*hex_to_rgb(cor_hex))
    badge.line.fill.background()
    add_text_box(slide7, 0.5, y, 1.6, 0.35, status, font_size=11, bold=True, color='#ffffff', align=PP_ALIGN.CENTER)

    add_text_box(slide7, 2.3, y - 0.05, 10, 0.35, titulo, font_size=16, bold=True, color=COR_PRIMARY)
    add_text_box(slide7, 2.3, y + 0.35, 10.5, 0.7, desc, font_size=11, color=COR_TEXT)
    y += 0.9

# ===== SLIDE 8: PRÓXIMOS PASSOS =====
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide8, COR_PRIMARY)
add_text_box(slide8, 1, 1, 11, 1, 'Próximos Passos', font_size=36, bold=True, color='#ffffff', align=PP_ALIGN.CENTER)

passos = [
    '1. Validar análise com equipe (Helaine, Cibele, Estética)',
    '2. Simular realocação dos 260 atendimentos de sábado na semana',
    '3. Definir cronograma de comunicação aos pacientes',
    '4. Período de transição: 30-60 dias',
    '5. Monitorar taxa de ocupação após mudança (métrica contínua)',
    f'6. Meta: implementar novo horário até Setembro/2026',
]
y = 2.5
for p in passos:
    add_text_box(slide8, 2, y, 9, 0.5, p, font_size=16, color='#ecf0f1')
    y += 0.55

# Linha final
shape = slide8.shapes.add_shape(1, Inches(4), Inches(6.5), Inches(5.33), Inches(0.03))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(*hex_to_rgb(COR_ACCENT))
shape.line.fill.background()

pptx_path = f'{BASE}/Relatorio_Taxa_Ociosidade_Relive.pptx'
prs.save(pptx_path)
print(f'[OK] PPTX salvo: {pptx_path}')

# ============================================================
# 2. GERAR PDF
# ============================================================
print("Gerando PDF...")
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import mm, cm
from reportlab.lib.colors import HexColor, white, black
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, Image
from reportlab.graphics.shapes import Drawing, Rect, String, Line
from reportlab.graphics.charts.barcharts import VerticalBarChart, HorizontalBarChart
from reportlab.graphics.charts.piecharts import Pie

pdf_path = f'{BASE}/Relatorio_Taxa_Ociosidade_Relive.pdf'
doc = SimpleDocTemplate(pdf_path, pagesize=A4,
                        leftMargin=20*mm, rightMargin=20*mm,
                        topMargin=20*mm, bottomMargin=20*mm)

styles = getSampleStyleSheet()
styleTitle = ParagraphStyle('CustomTitle', parent=styles['Title'], fontSize=22, textColor=HexColor(COR_PRIMARY), spaceAfter=6*mm)
styleH2 = ParagraphStyle('CustomH2', parent=styles['Heading2'], fontSize=16, textColor=HexColor(COR_PRIMARY), spaceAfter=4*mm, spaceBefore=8*mm)
styleH3 = ParagraphStyle('CustomH3', parent=styles['Heading3'], fontSize=13, textColor=HexColor(COR_SECONDARY), spaceAfter=3*mm, spaceBefore=6*mm)
styleBody = ParagraphStyle('CustomBody', parent=styles['Normal'], fontSize=10, textColor=HexColor(COR_TEXT), spaceAfter=3*mm, leading=14)
styleKpi = ParagraphStyle('KPI', parent=styles['Normal'], fontSize=24, textColor=HexColor(COR_PRIMARY), alignment=TA_CENTER, spaceAfter=0)
styleKpiLabel = ParagraphStyle('KPILabel', parent=styles['Normal'], fontSize=8, textColor=HexColor('#7f8c8d'), alignment=TA_CENTER, spaceAfter=2*mm)

elements = []

# CAPA
elements.append(Spacer(1, 40*mm))
elements.append(Paragraph('Relatório de Taxa de Ociosidade', styleTitle))
elements.append(Spacer(1, 5*mm))
elements.append(Paragraph('Análise de Viabilidade para Redução de Horário', styleH2))
elements.append(Paragraph('Abertura 8h · Fechamento 19h · Eliminação de Sábados', styleBody))
elements.append(Paragraph(f'Período: Janeiro a Junho 2026 | Relive', styleBody))
elements.append(Paragraph(f'Gerado em: {datetime.now().strftime("%d/%m/%Y %H:%M")}', styleBody))
elements.append(Spacer(1, 15*mm))

# Linha decorativa
d = Drawing(170*mm, 2*mm)
d.add(Line(0, 0, 170*mm, 0, strokeColor=HexColor(COR_ACCENT), strokeWidth=2))
elements.append(d)
elements.append(Spacer(1, 10*mm))

# RESUMO EXECUTIVO
elements.append(Paragraph('1. Resumo Executivo', styleH2))
elements.append(Paragraph(
    f'Este relatório analisa a viabilidade de três mudanças operacionais na clínica Relive: '
    f'<b>(a)</b> abertura às 8h em vez de 7h/7h30, <b>(b)</b> fechamento às 19h em vez de após 19h, '
    f'e <b>(c)</b> não abertura aos sábados. A análise compreende o período de 01/01/2026 a 30/06/2026.',
    styleBody
))

# Tabela de KPIs principais
kpi_table_data = [
    ['Indicador', 'Valor', 'Interpretação'],
    ['Total atendimentos impactados', str(kpis['total_geral']), 'Soma dos 3 cenários (6 meses)'],
    ['Atendimentos antes 8h', str(kpis['total_antes_8h']), f'Média {kpis["media_antes_8h_dia"]:.2f}/dia útil'],
    ['Atendimentos após 19h', str(kpis['total_apos_19h']), 'Média 4/mês'],
    ['Atendimentos sábados', str(kpis['total_sabados']), f'{kpis["sabados_unicos"]}/{kpis["total_sabados_periodo"]} sábados ocupados'],
    ['Taxa ocupação sábados', fmt_pct(kpis['taxa_ocupacao_sabados']), f'Média {kpis["media_sabado_dia"]:.1f} atend/sábado'],
    ['Pacientes únicos impactados', str(kpis['sabados_pacientes_unicos'] + kpis['antes_8h_pacientes_unicos'] + kpis['apos_19h_pacientes_unicos']), 'Todos os cenários'],
]

kpi_table = Table(kpi_table_data, colWidths=[60*mm, 35*mm, 65*mm])
kpi_table.setStyle(TableStyle([
    ('BACKGROUND', (0, 0), (-1, 0), HexColor(COR_PRIMARY)),
    ('TEXTCOLOR', (0, 0), (-1, 0), white),
    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
    ('FONTSIZE', (0, 0), (-1, 0), 10),
    ('FONTSIZE', (0, 1), (-1, -1), 9),
    ('ALIGN', (1, 0), (1, -1), 'CENTER'),
    ('GRID', (0, 0), (-1, -1), 0.5, HexColor('#bdc3c7')),
    ('ROWBACKGROUNDS', (0, 1), (-1, -1), [white, HexColor('#f8f9fa')]),
    ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
    ('TOPPADDING', (0, 0), (-1, -1), 4),
    ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
]))
elements.append(kpi_table)
elements.append(Spacer(1, 6*mm))

# Conclusão preliminar
elements.append(Paragraph('2. Conclusão Preliminar', styleH2))
conclusoes = [
    '<b>✓ Abertura às 8h: VIÁVEL</b> — Apenas 0,51 atendimentos/dia antes das 8h. Impacto mínimo.',
    '<b>✓ Fechamento às 19h: VIÁVEL</b> — Apenas 24 atendimentos após 19h em 6 meses.',
    '<b>⚠ Sábados: REQUER ATENÇÃO</b> — 80,8% de ocupação, 260 atendimentos. Eliminar exige plano.',
]
for c in conclusoes:
    elements.append(Paragraph(c, styleBody))

elements.append(PageBreak())

# ANÁLISE SÁBADOS
elements.append(Paragraph('3. Análise Detalhada: Sábados', styleH2))

# Gráfico mensal sábados
drawing = Drawing(170*mm, 80*mm)
bc = VerticalBarChart()
bc.x = 40
bc.y = 40
bc.height = 120
bc.width = 400
meses_sab = list(kpis['sabados_por_mes'].keys())
valores_sab = list(kpis['sabados_por_mes'].values())
bc.data = [valores_sab]
bc.categoryAxis.categoryNames = [m[-2:] for m in meses_sab]  # só mês
bc.categoryAxis.labels.fontSize = 8
bc.valueAxis.valueMin = 0
bc.valueAxis.valueMax = max(valores_sab) * 1.2
bc.valueAxis.valueStep = 10
bc.bars[0].fillColor = HexColor(COR_ACCENT)
bc.barWidth = 15
bc.barSpacing = 8
drawing.add(bc)
drawing.add(String(200, 175, 'Atendimentos aos Sábados por Mês (2026)', fontSize=11, fillColor=HexColor(COR_PRIMARY), textAnchor='middle'))
elements.append(drawing)
elements.append(Spacer(1, 5*mm))

# Tabela sábados
sab_table_data = [['Métrica', 'Valor']]
sab_metrics = [
    ('Total atendimentos', str(kpis['total_sabados'])),
    ('Sábados com atendimento', f'{kpis["sabados_unicos"]}/{kpis["total_sabados_periodo"]}'),
    ('Sábados sem atendimento', str(kpis['sabados_sem_atendimento'])),
    ('Taxa de ocupação', fmt_pct(kpis['taxa_ocupacao_sabados'])),
    ('Média atendimentos/sábado', f'{kpis["media_sabado_dia"]:.1f}'),
    ('Pacientes únicos', str(kpis['sabados_pacientes_unicos'])),
]
for m, v in sab_metrics:
    sab_table_data.append([m, v])

sab_table = Table(sab_table_data, colWidths=[80*mm, 60*mm])
sab_table.setStyle(TableStyle([
    ('BACKGROUND', (0, 0), (-1, 0), HexColor(COR_ACCENT)),
    ('TEXTCOLOR', (0, 0), (-1, 0), white),
    ('GRID', (0, 0), (-1, -1), 0.5, HexColor('#bdc3c7')),
    ('ROWBACKGROUNDS', (0, 1), (-1, -1), [white, HexColor('#f8f9fa')]),
    ('FONTSIZE', (0, 0), (-1, -1), 9),
    ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
    ('TOPPADDING', (0, 0), (-1, -1), 3),
    ('BOTTOMPADDING', (0, 0), (-1, -1), 3),
]))
elements.append(sab_table)
elements.append(Spacer(1, 5*mm))

# Top tipos sábado
elements.append(Paragraph('Top 5 Tipos de Atendimento (Sábados):', styleH3))
for tipo, qtd in list(kpis['sabados_por_tipo'].items())[:5]:
    elements.append(Paragraph(f'• {tipo}: <b>{qtd}</b> atendimentos', styleBody))

elements.append(PageBreak())

# ANTES 8H
elements.append(Paragraph('4. Análise Detalhada: Antes das 8h', styleH2))

drawing2 = Drawing(170*mm, 80*mm)
bc2 = VerticalBarChart()
bc2.x = 40
bc2.y = 40
bc2.height = 120
bc2.width = 400
meses_8h = list(kpis['antes_8h_por_mes'].keys())
valores_8h = list(kpis['antes_8h_por_mes'].values())
bc2.data = [valores_8h]
bc2.categoryAxis.categoryNames = [m[-2:] for m in meses_8h]
bc2.categoryAxis.labels.fontSize = 8
bc2.valueAxis.valueMin = 0
bc2.valueAxis.valueMax = max(valores_8h) * 1.3 if valores_8h else 20
bc2.valueAxis.valueStep = 5
bc2.bars[0].fillColor = HexColor(COR_ORANGE)
bc2.barWidth = 15
bc2.barSpacing = 8
drawing2.add(bc2)
drawing2.add(String(200, 175, 'Atendimentos Antes 8h por Mês (2026)', fontSize=11, fillColor=HexColor(COR_PRIMARY), textAnchor='middle'))
elements.append(drawing2)
elements.append(Spacer(1, 5*mm))

elements.append(Paragraph(
    f'Total de <b>{kpis["total_antes_8h"]}</b> atendimentos antes das 8h em {kpis["dias_uteis"]} dias úteis. '
    f'Média de <b>{kpis["media_antes_8h_dia"]:.2f}</b> atendimentos por dia útil. '
    f'Profissionais envolvidos: Helaine Beatriz Jacobucci (42) e Cibele Priscila Busch Furlan (22).',
    styleBody
))

elements.append(Paragraph('Tipos de atendimento antes 8h:', styleH3))
for tipo, qtd in kpis['antes_8h_por_tipo'].items():
    elements.append(Paragraph(f'• {tipo}: <b>{qtd}</b>', styleBody))

elements.append(PageBreak())

# APÓS 19H
elements.append(Paragraph('5. Análise Detalhada: Após 19h', styleH2))

drawing3 = Drawing(170*mm, 70*mm)
bc3 = HorizontalBarChart()
bc3.x = 50
bc3.y = 10
bc3.height = 100
bc3.width = 350
tipos_19h = [t[:40] for t in list(kpis['apos_19h_por_tipo'].keys())]
valores_19h = list(kpis['apos_19h_por_tipo'].values())
bc3.data = [valores_19h]
bc3.categoryAxis.categoryNames = tipos_19h
bc3.categoryAxis.labels.fontSize = 7
bc3.valueAxis.valueMin = 0
bc3.valueAxis.valueMax = max(valores_19h) + 2 if valores_19h else 8
bc3.valueAxis.valueStep = 1
bc3.bars[0].fillColor = HexColor(COR_GREEN)
drawing3.add(bc3)
drawing3.add(String(200, 120, 'Tipos de Atendimento Após 19h', fontSize=10, fillColor=HexColor(COR_PRIMARY), textAnchor='middle'))
elements.append(drawing3)
elements.append(Spacer(1, 5*mm))

elements.append(Paragraph(
    f'Apenas <b>{kpis["total_apos_19h"]}</b> atendimentos após 19h em 6 meses. '
    f'<b>{kpis["apos_19h_por_profissional"].get("Estética", 0) + kpis["apos_19h_por_profissional"].get("Estética Grazi", 0) + kpis["apos_19h_por_profissional"].get("Estética Lais", 0)}</b> '
    f'são da área de estética, representando a maioria absoluta.',
    styleBody
))

elements.append(PageBreak())

# IMPACTO POR PROFISSIONAL
elements.append(Paragraph('6. Impacto por Profissional (Todos os Cenários)', styleH2))

prof_table_data = [['Profissional', 'Atendimentos', '% do Total']]
for prof, qtd in kpis['top10_profissionais_impactados'].items():
    pct = qtd / kpis['total_geral'] * 100
    prof_table_data.append([prof, str(qtd), fmt_pct(pct)])

prof_table = Table(prof_table_data, colWidths=[80*mm, 40*mm, 40*mm])
prof_table.setStyle(TableStyle([
    ('BACKGROUND', (0, 0), (-1, 0), HexColor(COR_PRIMARY)),
    ('TEXTCOLOR', (0, 0), (-1, 0), white),
    ('GRID', (0, 0), (-1, -1), 0.5, HexColor('#bdc3c7')),
    ('ROWBACKGROUNDS', (0, 1), (-1, -1), [white, HexColor('#f8f9fa')]),
    ('FONTSIZE', (0, 0), (-1, -1), 9),
    ('ALIGN', (1, 0), (-1, -1), 'CENTER'),
    ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
    ('TOPPADDING', (0, 0), (-1, -1), 4),
    ('BOTTOMPADDING', (0, 0), (-1, -1), 4),
]))
elements.append(prof_table)
elements.append(Spacer(1, 10*mm))

# RECOMENDAÇÕES
elements.append(Paragraph('7. Recomendações', styleH2))

recs = [
    ('<b>Abertura às 8h:</b> VIÁVEL. Implementar imediatamente. Rearranjar horários de Helaine e Cibele para 8h. '
     'Comunicar pacientes afetados com 15 dias de antecedência.'),
    ('<b>Fechamento às 19h:</b> VIÁVEL. Implementar com transição gradual. Avaliar possibilidade de manter '
     '1 dia/semana com horário estendido para estética (ex: quartas até 20h).'),
    ('<b>Sábados:</b> NÃO ELIMINAR completamente. Recomendação: reduzir para 2 sábados/mês '
     '(1º e 3º sábado). Isso reduziria ~50% dos atendimentos de sábado mantendo disponibilidade. '
     'Reavaliar após 6 meses.'),
    ('<b>Plano de transição:</b> 30-60 dias. Comunicar pacientes, reorganizar agenda, treinar equipe. '
     'Meta: novo horário operacional até Setembro/2026.'),
]
for r in recs:
    elements.append(Paragraph(r, styleBody))
    elements.append(Spacer(1, 3*mm))

elements.append(Spacer(1, 10*mm))
d2 = Drawing(170*mm, 2*mm)
d2.add(Line(0, 0, 170*mm, 0, strokeColor=HexColor(COR_ACCENT), strokeWidth=1))
elements.append(d2)
elements.append(Paragraph('Relatório gerado automaticamente — Dados: Janeiro a Junho 2026', styleKpiLabel))

doc.build(elements)
print(f'[OK] PDF salvo: {pdf_path}')

# ============================================================
# 3. GERAR HTML
# ============================================================
print("Gerando HTML...")

html_content = f'''<!DOCTYPE html>
<html lang="pt-BR">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>Relatório de Taxa de Ociosidade — Relive</title>
<script src="https://cdn.jsdelivr.net/npm/chart.js@4.4.0/dist/chart.umd.min.js"></script>
<style>
* {{ margin: 0; padding: 0; box-sizing: border-box; }}
body {{ font-family: 'Segoe UI', system-ui, -apple-system, sans-serif; background: #f0f2f5; color: #2c3e50; line-height: 1.6; }}
.header {{ background: linear-gradient(135deg, {COR_PRIMARY}, {COR_SECONDARY}); color: white; padding: 60px 40px; text-align: center; }}
.header h1 {{ font-size: 2.5em; margin-bottom: 10px; }}
.header p {{ font-size: 1.2em; opacity: 0.85; }}
.container {{ max-width: 1300px; margin: 0 auto; padding: 30px 20px; }}
.kpi-grid {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(280px, 1fr)); gap: 20px; margin-bottom: 40px; }}
.kpi-card {{ background: white; border-radius: 12px; padding: 25px; box-shadow: 0 2px 8px rgba(0,0,0,0.08); border-top: 4px solid {COR_PRIMARY}; }}
.kpi-card.warning {{ border-top-color: {COR_ACCENT}; }}
.kpi-card.success {{ border-top-color: {COR_GREEN}; }}
.kpi-card.info {{ border-top-color: {COR_ORANGE}; }}
.kpi-card .kpi-label {{ font-size: 0.8em; color: #7f8c8d; text-transform: uppercase; letter-spacing: 1px; margin-bottom: 8px; }}
.kpi-card .kpi-value {{ font-size: 2.2em; font-weight: 700; color: {COR_PRIMARY}; }}
.kpi-card.warning .kpi-value {{ color: {COR_ACCENT}; }}
.kpi-card.success .kpi-value {{ color: {COR_GREEN}; }}
.kpi-card.info .kpi-value {{ color: {COR_ORANGE}; }}
.kpi-card .kpi-sub {{ font-size: 0.85em; color: #95a5a6; margin-top: 5px; }}
.section {{ background: white; border-radius: 12px; padding: 30px; margin-bottom: 30px; box-shadow: 0 2px 8px rgba(0,0,0,0.08); }}
.section h2 {{ font-size: 1.6em; color: {COR_PRIMARY}; margin-bottom: 15px; padding-bottom: 10px; border-bottom: 2px solid #ecf0f1; }}
.section h3 {{ font-size: 1.2em; color: {COR_SECONDARY}; margin: 20px 0 10px; }}
.chart-row {{ display: grid; grid-template-columns: 1fr 1fr; gap: 30px; margin: 20px 0; }}
.chart-box {{ background: #fafafa; border-radius: 8px; padding: 20px; }}
.chart-box.full {{ grid-column: 1 / -1; }}
.chart-box canvas {{ max-height: 350px; }}
.conclusion-box {{ background: #eaf2f8; border-left: 5px solid {COR_PRIMARY}; padding: 20px 25px; border-radius: 0 8px 8px 0; margin: 20px 0; }}
.conclusion-box.warning {{ background: #fef9e7; border-left-color: {COR_ORANGE}; }}
.conclusion-box.success {{ background: #eafaf1; border-left-color: {COR_GREEN}; }}
.conclusion-box h4 {{ margin-bottom: 8px; }}
table {{ width: 100%; border-collapse: collapse; margin: 15px 0; }}
th {{ background: {COR_PRIMARY}; color: white; padding: 10px 15px; text-align: left; font-size: 0.9em; }}
td {{ padding: 8px 15px; border-bottom: 1px solid #ecf0f1; font-size: 0.9em; }}
tr:nth-child(even) {{ background: #f8f9fa; }}
.badge {{ display: inline-block; padding: 4px 12px; border-radius: 20px; font-size: 0.75em; font-weight: 700; text-transform: uppercase; }}
.badge-ok {{ background: #d5f5e3; color: #1e8449; }}
.badge-warn {{ background: #fdebd0; color: #b9770e; }}
.badge-alert {{ background: #fadbd8; color: #c0392b; }}
.footer {{ text-align: center; padding: 30px; color: #95a5a6; font-size: 0.85em; }}
.progress-bar {{ height: 8px; background: #ecf0f1; border-radius: 4px; margin-top: 5px; overflow: hidden; }}
.progress-fill {{ height: 100%; border-radius: 4px; }}
@media (max-width: 768px) {{ .chart-row {{ grid-template-columns: 1fr; }} .kpi-grid {{ grid-template-columns: 1fr; }} }}
</style>
</head>
<body>

<div class="header">
    <h1>📊 Relatório de Taxa de Ociosidade</h1>
    <p>Análise de Viabilidade: Abertura 8h · Fechamento 19h · Sábados</p>
    <p style="margin-top:10px;font-size:0.95em;">Período: Janeiro a Junho 2026 | Relive</p>
</div>

<div class="container">

    <!-- KPIs Principais -->
    <div class="kpi-grid">
        <div class="kpi-card warning">
            <div class="kpi-label">Total Atendimentos Impactados</div>
            <div class="kpi-value">{kpis['total_geral']}</div>
            <div class="kpi-sub">Soma dos 3 cenários (6 meses)</div>
        </div>
        <div class="kpi-card">
            <div class="kpi-label">Sábados com Atendimento</div>
            <div class="kpi-value">{kpis['sabados_unicos']}/{kpis['total_sabados_periodo']}</div>
            <div class="kpi-sub">Taxa de ocupação: {fmt_pct(kpis['taxa_ocupacao_sabados'])}</div>
            <div class="progress-bar"><div class="progress-fill" style="width:{kpis['taxa_ocupacao_sabados']}%;background:{COR_ACCENT};"></div></div>
        </div>
        <div class="kpi-card info">
            <div class="kpi-label">Média/Dia Antes 8h</div>
            <div class="kpi-value">{kpis['media_antes_8h_dia']:.2f}</div>
            <div class="kpi-sub">{kpis['total_antes_8h']} atendimentos em {kpis['dias_uteis']} dias úteis</div>
        </div>
        <div class="kpi-card success">
            <div class="kpi-label">Atendimentos Após 19h</div>
            <div class="kpi-value">{kpis['total_apos_19h']}</div>
            <div class="kpi-sub">Em 6 meses de análise</div>
        </div>
    </div>

    <!-- Conclusão -->
    <div class="section">
        <h2>Conclusão Preliminar</h2>
        <div class="conclusion-box success">
            <h4>✓ Abertura às 8h: VIÁVEL</h4>
            <p>Apenas 0,51 atendimentos/dia antes das 8h (66 em 129 dias úteis). Impacto mínimo na operação, concentrado em nutrição (Helaine e Cibele).</p>
        </div>
        <div class="conclusion-box success">
            <h4>✓ Fechamento às 19h: VIÁVEL</h4>
            <p>Apenas 24 atendimentos após 19h em 6 meses. Maioria de estética (16 de 24), possível reacomodar dentro do horário comercial.</p>
        </div>
        <div class="conclusion-box warning">
            <h4>⚠ Sábados: REQUER ANÁLISE CUIDADOSA</h4>
            <p>80,8% dos sábados têm atendimento, média de 12,4 atendimentos/dia, 260 no total. Eliminar exigiria redistribuir aproximadamente 43 atendimentos/mês.</p>
        </div>
        <div class="conclusion-box" style="background:#fdedec;border-left-color:{COR_ACCENT};">
            <h4>⚠ Cenário Combinado: REQUER PLANO DE TRANSIÇÃO</h4>
            <p>350 atendimentos impactados em 6 meses (~58/mês). Necessário plano estruturado: migrar sábados para dias úteis, reacomodar estética noturna, comunicar pacientes com antecedência.</p>
        </div>
    </div>

    <!-- Sábados -->
    <div class="section">
        <h2>📅 Análise: Sábados</h2>
        <div class="kpi-grid" style="margin-bottom:20px;">
            <div class="kpi-card warning">
                <div class="kpi-label">Total Atendimentos</div>
                <div class="kpi-value">{kpis['total_sabados']}</div>
            </div>
            <div class="kpi-card warning">
                <div class="kpi-label">Média por Sábado</div>
                <div class="kpi-value">{kpis['media_sabado_dia']:.1f}</div>
            </div>
            <div class="kpi-card warning">
                <div class="kpi-label">Sábados sem Atendimento</div>
                <div class="kpi-value">{kpis['sabados_sem_atendimento']}/{kpis['total_sabados_periodo']}</div>
            </div>
            <div class="kpi-card warning">
                <div class="kpi-label">Pacientes Únicos</div>
                <div class="kpi-value">{kpis['sabados_pacientes_unicos']}</div>
            </div>
        </div>
        <div class="chart-row">
            <div class="chart-box">
                <h3>Atendimentos por Mês</h3>
                <canvas id="chartSabadosMes"></canvas>
            </div>
            <div class="chart-box">
                <h3>Top 8 Tipos de Atendimento</h3>
                <canvas id="chartSabadosTipo"></canvas>
            </div>
        </div>
    </div>

    <!-- Antes 8h -->
    <div class="section">
        <h2>🌅 Análise: Antes das 8h</h2>
        <div class="kpi-grid" style="margin-bottom:20px;">
            <div class="kpi-card info">
                <div class="kpi-label">Total</div>
                <div class="kpi-value">{kpis['total_antes_8h']}</div>
            </div>
            <div class="kpi-card info">
                <div class="kpi-label">Média por Dia Útil</div>
                <div class="kpi-value">{kpis['media_antes_8h_dia']:.2f}</div>
            </div>
            <div class="kpi-card info">
                <div class="kpi-label">Dias Úteis no Período</div>
                <div class="kpi-value">{kpis['dias_uteis']}</div>
            </div>
            <div class="kpi-card info">
                <div class="kpi-label">Pacientes Únicos</div>
                <div class="kpi-value">{kpis['antes_8h_pacientes_unicos']}</div>
            </div>
        </div>
        <div class="chart-row">
            <div class="chart-box full">
                <h3>Atendimentos por Mês (Antes 8h)</h3>
                <canvas id="chartAntes8hMes"></canvas>
            </div>
        </div>
    </div>

    <!-- Após 19h -->
    <div class="section">
        <h2>🌙 Análise: Após 19h</h2>
        <div class="kpi-grid" style="margin-bottom:20px;">
            <div class="kpi-card success">
                <div class="kpi-label">Total</div>
                <div class="kpi-value">{kpis['total_apos_19h']}</div>
            </div>
            <div class="kpi-card success">
                <div class="kpi-label">Pacientes Únicos</div>
                <div class="kpi-value">{kpis['apos_19h_pacientes_unicos']}</div>
            </div>
            <div class="kpi-card success">
                <div class="kpi-label">Principal Área</div>
                <div class="kpi-value" style="font-size:1.4em;">Estética</div>
                <div class="kpi-sub">16/24 atendimentos (66,7%)</div>
            </div>
        </div>
        <div class="chart-box">
            <h3>Tipos de Atendimento (Após 19h)</h3>
            <canvas id="chartApos19hTipo"></canvas>
        </div>
    </div>

    <!-- Impacto por Profissional -->
    <div class="section">
        <h2>👩‍⚕️ Impacto por Profissional</h2>
        <div class="chart-row">
            <div class="chart-box full">
                <canvas id="chartProfissionais"></canvas>
            </div>
        </div>
        <table>
            <tr><th>Profissional</th><th>Atendimentos</th><th>% do Total</th><th>Avaliação</th></tr>
'''

for prof, qtd in kpis['top10_profissionais_impactados'].items():
    pct = qtd / kpis['total_geral'] * 100
    badge_class = 'badge-alert' if pct > 30 else ('badge-warn' if pct > 10 else 'badge-ok')
    html_content += f'<tr><td>{prof}</td><td><b>{qtd}</b></td><td>{fmt_pct(pct)}</td><td><span class="badge {badge_class}">{"ALTO" if pct > 30 else ("MÉDIO" if pct > 10 else "BAIXO")}</span></td></tr>\n'

html_content += f'''
        </table>
    </div>

    <!-- Recomendações -->
    <div class="section">
        <h2>📋 Recomendações Finais</h2>
        <table>
            <tr><th>Cenário</th><th>Viabilidade</th><th>Recomendação</th></tr>
            <tr>
                <td><b>Abertura às 8h</b></td>
                <td><span class="badge badge-ok">VIÁVEL</span></td>
                <td>Implementar imediatamente. Rearranjar horários de Helaine (42) e Cibele (22) para 8h. Comunicar pacientes com 15 dias.</td>
            </tr>
            <tr>
                <td><b>Fechamento às 19h</b></td>
                <td><span class="badge badge-ok">VIÁVEL</span></td>
                <td>Implementar com transição gradual. Avaliar 1 dia/semana com horário estendido para estética.</td>
            </tr>
            <tr>
                <td><b>Eliminar Sábados</b></td>
                <td><span class="badge badge-warn">ATENÇÃO</span></td>
                <td>Não eliminar completamente. Reduzir para 2 sábados/mês (1º e 3º). Reavaliar após 6 meses.</td>
            </tr>
            <tr>
                <td><b>Cenário Combinado</b></td>
                <td><span class="badge badge-alert">PLANO NECESSÁRIO</span></td>
                <td>350 atendimentos impactados. Plano de transição 30-60 dias. Meta: Setembro/2026.</td>
            </tr>
        </table>
    </div>

    <!-- Distribuição por Hora (Sábados) -->
    <div class="section">
        <h2>⏰ Distribuição por Horário (Sábados)</h2>
        <div class="chart-box">
            <canvas id="chartHorasSabado"></canvas>
        </div>
    </div>

    <div class="footer">
        <p>Relatório gerado em {datetime.now().strftime("%d/%m/%Y %H:%M")} | Dados: Janeiro a Junho 2026 | Relive</p>
        <p>Análise de Taxa de Ociosidade para decisão operacional</p>
    </div>
</div>

<script>
// Cores
const primary = '{COR_PRIMARY}';
const secondary = '{COR_SECONDARY}';
const accent = '{COR_ACCENT}';
const green = '{COR_GREEN}';
const orange = '{COR_ORANGE}';

// Chart.js defaults
Chart.defaults.font.family = "'Segoe UI', system-ui, sans-serif";
Chart.defaults.font.size = 12;

// 1. Sábados por Mês
new Chart(document.getElementById('chartSabadosMes'), {{
    type: 'bar',
    data: {{
        labels: {json.dumps([m[-2:] + '/' + m[:4] for m in kpis['sabados_por_mes'].keys()])},
        datasets: [{{
            label: 'Atendimentos',
            data: {json.dumps(list(kpis['sabados_por_mes'].values()))},
            backgroundColor: accent,
            borderRadius: 6,
        }}]
    }},
    options: {{
        responsive: true,
        plugins: {{ legend: {{ display: false }} }},
        scales: {{ y: {{ beginAtZero: true, ticks: {{ stepSize: 10 }} }} }}
    }}
}});

// 2. Sábados por Tipo
new Chart(document.getElementById('chartSabadosTipo'), {{
    type: 'bar',
    data: {{
        labels: {json.dumps([t[:30] for t in list(kpis['sabados_por_tipo'].keys())[:8]])},
        datasets: [{{
            label: 'Atendimentos',
            data: {json.dumps(list(kpis['sabados_por_tipo'].values())[:8])},
            backgroundColor: accent + 'cc',
            borderRadius: 6,
        }}]
    }},
    options: {{
        indexAxis: 'y',
        responsive: true,
        plugins: {{ legend: {{ display: false }} }},
        scales: {{ x: {{ beginAtZero: true, ticks: {{ stepSize: 5 }} }} }}
    }}
}});

// 3. Antes 8h por Mês
new Chart(document.getElementById('chartAntes8hMes'), {{
    type: 'bar',
    data: {{
        labels: {json.dumps([m[-2:] + '/' + m[:4] for m in kpis['antes_8h_por_mes'].keys()])},
        datasets: [{{
            label: 'Atendimentos',
            data: {json.dumps(list(kpis['antes_8h_por_mes'].values()))},
            backgroundColor: orange,
            borderRadius: 6,
        }}]
    }},
    options: {{
        responsive: true,
        plugins: {{ legend: {{ display: false }} }},
        scales: {{ y: {{ beginAtZero: true, ticks: {{ stepSize: 5 }} }} }}
    }}
}});

// 4. Após 19h por Tipo
new Chart(document.getElementById('chartApos19hTipo'), {{
    type: 'bar',
    data: {{
        labels: {json.dumps([t[:30] for t in kpis['apos_19h_por_tipo'].keys()])},
        datasets: [{{
            label: 'Atendimentos',
            data: {json.dumps(list(kpis['apos_19h_por_tipo'].values()))},
            backgroundColor: green,
            borderRadius: 6,
        }}]
    }},
    options: {{
        indexAxis: 'y',
        responsive: true,
        plugins: {{ legend: {{ display: false }} }},
        scales: {{ x: {{ beginAtZero: true, ticks: {{ stepSize: 1 }} }} }}
    }}
}});

// 5. Profissionais Impactados
new Chart(document.getElementById('chartProfissionais'), {{
    type: 'bar',
    data: {{
        labels: {json.dumps([p[:30] for p in kpis['top10_profissionais_impactados'].keys()])},
        datasets: [{{
            label: 'Atendimentos Impactados',
            data: {json.dumps(list(kpis['top10_profissionais_impactados'].values()))},
            backgroundColor: [accent, accent, orange, orange, green, green, secondary, secondary],
            borderRadius: 6,
        }}]
    }},
    options: {{
        responsive: true,
        plugins: {{ legend: {{ display: false }} }},
        scales: {{ y: {{ beginAtZero: true }} }}
    }}
}});

// 6. Distribuição Horária Sábados
new Chart(document.getElementById('chartHorasSabado'), {{
    type: 'line',
    data: {{
        labels: {json.dumps(list(kpis['sabados_por_hora'].keys()))},
        datasets: [{{
            label: 'Atendimentos por Horário',
            data: {json.dumps(list(kpis['sabados_por_hora'].values()))},
            borderColor: accent,
            backgroundColor: accent + '33',
            fill: true,
            tension: 0.3,
            pointRadius: 5,
            pointBackgroundColor: accent,
        }}]
    }},
    options: {{
        responsive: true,
        plugins: {{ legend: {{ display: false }} }},
        scales: {{ y: {{ beginAtZero: true }} }}
    }}
}});
</script>

</body>
</html>
'''

html_path = f'{BASE}/Relatorio_Taxa_Ociosidade_Relive.html'
with open(html_path, 'w', encoding='utf-8') as f:
    f.write(html_content)
print(f'[OK] HTML salvo: {html_path}')

print('\n=== TODOS OS RELATORIOS GERADOS COM SUCESSO ===')
print(f'[PPTX] {pptx_path}')
print(f'[PDF]  {pdf_path}')
print(f'[HTML] {html_path}')
